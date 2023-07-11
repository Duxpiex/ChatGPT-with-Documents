import os
import tempfile
from io import BufferedReader
from typing import Optional
from fastapi import UploadFile
import mimetypes
from PyPDF2 import PdfReader
import docx2txt
import csv
import pptx
from loguru import logger

from models.models import Document, DocumentMetadata


async def get_document_from_file(
    file: UploadFile, metadata: DocumentMetadata
) -> Document:
    # 임시 파일 생성
    with tempfile.NamedTemporaryFile(delete=False) as temp_file:
        # 파일에 쓰기
        content = await file.read()
        temp_file.write(content)
        temp_file_path = temp_file.name

    # 임시 파일의 경로를 사용하여 텍스트 추출
    extracted_text = extract_text_from_filepath(temp_file_path, file.content_type)

    doc = Document(text=extracted_text, metadata=metadata)

    # 임시 파일 삭제
    os.remove(temp_file_path)

    return doc


def extract_text_from_filepath(filepath: str, mimetype: Optional[str] = None) -> str:
    """Return the text content of a file given its filepath."""

    if mimetype is None:
        # Get the mimetype of the file based on its extension
        mimetype, _ = mimetypes.guess_type(filepath)

    if not mimetype:
        if filepath.endswith(".md"):
            mimetype = "text/markdown"
        else:
            raise Exception("Unsupported file type")

    try:
        # For some types of files, we need to pass the path instead of the file object
        if mimetype in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
            extracted_text = extract_text_from_file(filepath, mimetype)
        else:
            with open(filepath, "rb") as file:
                extracted_text = extract_text_from_file(file, mimetype)
    except Exception as e:
        logger.error(e)
        raise e

    return extracted_text


def extract_text_from_file(file, mimetype: str) -> str:
    if mimetype == "application/pdf":
        # Extract text from pdf using PyPDF2
        reader = PdfReader(file)
        extracted_text = " ".join([page.extract_text() for page in reader.pages])
    elif mimetype == "text/plain" or mimetype == "text/markdown":
        # Read text from plain text file
        extracted_text = file.read().decode("utf-8")
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        # Extract text from docx using docx2txt
        extracted_text = docx2txt.process(file)
    elif mimetype == "text/csv":
        # Extract text from csv using csv module
        extracted_text = ""
        decoded_buffer = (line.decode("utf-8") for line in file)
        reader = csv.reader(decoded_buffer)
        for row in reader:
            extracted_text += " ".join(row) + "\n"
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        # Extract text from pptx using python-pptx
        extracted_text = ""
        presentation = pptx.Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                    extracted_text += "\n"
    else:
        # Unsupported file type
        raise ValueError("Unsupported file type: {}".format(mimetype))

    return extracted_text
